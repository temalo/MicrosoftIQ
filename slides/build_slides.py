#!/usr/bin/env python3
"""
Build the Microsoft IQ pitch deck for the Contoso Events demo.

    python build_slides.py            # -> microsoft-iq-story.pptx (5 slides)

Numbers are read from ../data/output/manifest.json, so the deck always matches
the generated data. Run `python ../generate.py` first. Brand via BRAND env var.
Requires python-pptx (see ../requirements.txt).
"""

import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BRAND = os.environ.get("BRAND", "Contoso Events")
HERE = os.path.dirname(os.path.abspath(__file__))
MANIFEST = os.path.join(HERE, "..", "data", "output", "manifest.json")
OUT = os.path.join(HERE, "microsoft-iq-story.pptx")

NAVY = RGBColor(0x1B, 0x2A, 0x4A); BLUE = RGBColor(0x2F, 0x6B, 0xD6)
BLUED = RGBColor(0x24, 0x56, 0xB3); INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x5B, 0x67, 0x70); LINE = RGBColor(0xE2, 0xE6, 0xEA)
BG = RGBColor(0xF4, 0xF6, 0xF8); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE = RGBColor(0xE9, 0xF0, 0xFB); GREEN = RGBColor(0x1E, 0x7A, 0x3D)
AMBER = RGBColor(0xB0, 0x6A, 0x00); CFE0 = RGBColor(0xCF, 0xE0, 0xF2)

# --- load headline numbers --------------------------------------------------
if os.path.exists(MANIFEST):
    M = json.load(open(MANIFEST, encoding="utf-8"))
else:
    print("manifest.json not found - run generate.py first. Using placeholders.")
    M = {"top_speakers_by_licensed_users_attended": [["Speaker One", 4000],
         ["Speaker Two", 3600], ["Speaker Three", 2500], ["Speaker Four", 2400],
         ["Speaker Five", 2300]],
         "top_sponsors_by_influenced_pipeline": [["Sponsor A", 5200000, 2.6],
         ["Sponsor B", 1800000, 2.4], ["Sponsor C", 1150000, 2.3],
         ["Sponsor D", 1010000, 2.2], ["Sponsor E", 980000, 2.3]],
         "licensed_user_definitions": {"NamedLicensedUser": 658,
         "PaidLicensedSeat": 749, "ServedPopulation": 2294,
         "LoginCapableIdentity": 3484}}

TOP_SPK = M["top_speakers_by_licensed_users_attended"]
TOP_SPN = M["top_sponsors_by_influenced_pipeline"]
DEFS = M["licensed_user_definitions"]

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _noline(sh): sh.line.fill.background()


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: _noline(shp)
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    if isinstance(runs[0], tuple): runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Segoe UI"
            if rest and rest[0] == "i": r.font.italic = True
    return tb


def header(slide, tag="ILLUSTRATIVE"):
    rect(slide, 0, 0, SW, Inches(0.92), fill=NAVY)
    rect(slide, 0, Inches(0.92), SW, Pt(3), fill=BLUE)
    text(slide, Inches(0.55), Inches(0.16), Inches(9), Inches(0.6),
         [[(BRAND, 22, WHITE, True), ("    |    ", 18, RGBColor(0x6f,0x86,0xa6), False),
           ("Conferences Intelligence", 14, CFE0, True)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, SW - Inches(2.35), Inches(0.26), Inches(1.8), Inches(0.4),
         fill=None, line=RGBColor(0x6f,0x86,0xa6), line_w=1.0, round_=True)
    text(slide, SW - Inches(2.35), Inches(0.28), Inches(1.8), Inches(0.36),
         [[("\u25CF ", 9, BLUE, True), (tag, 10.5, CFE0, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    rect(slide, 0, SH - Inches(0.52), SW, Inches(0.52), fill=RGBColor(0xED,0xF1,0xF5))
    text(slide, Inches(0.55), SH - Inches(0.5), Inches(12.2), Inches(0.46),
         [[("Illustrative only. ", 9, NAVY, True),
           (f"Synthetic demo data generated for this project. \u201c{BRAND}\u201d is a "
            "fictitious company; any resemblance to real organizations is coincidental.",
            9, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)


def bg(slide): rect(slide, 0, 0, SW, SH, fill=BG)


def fmt_m(v): return f"${v/1_000_000:.1f}M"


# ---------------- SLIDE 1 : the output --------------------------------------
s = prs.slides.add_slide(BLANK); bg(s); header(s); footer(s)
text(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.6),
     [[("What the output of this exercise looks like", 27, NAVY, True)]])
text(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.5),
     [[("One business question \u2014 ", 16, GREY, False),
       ("\u201CHow many licensed users do we have?\u201D", 16, BLUE, True),
       (" \u2014 four defensible, reconcilable answers.", 16, GREY, False)]])
rect(s, Inches(0.55), Inches(2.5), Inches(6.0), Inches(0.72), fill=LBLUE, round_=True)
text(s, Inches(0.85), Inches(2.5), Inches(5.5), Inches(0.72),
     [[("You: ", 13, GREY, True), ("How many licensed users do we have?", 13, INK, False)]],
     anchor=MSO_ANCHOR.MIDDLE)
tx, ty, tw = Inches(0.55), Inches(3.4), Inches(7.5); rowh = Inches(0.52)
rect(s, tx, ty, tw, Inches(0.5), fill=NAVY)
cols = [Inches(2.5), Inches(3.1), Inches(1.0), Inches(0.9)]; cxs = [tx]
for c in cols[:-1]: cxs.append(cxs[-1] + c)
for i, (h, cw) in enumerate(zip(["Business unit", "Definition", "Count", "Status"], cols)):
    al = PP_ALIGN.LEFT if i < 2 else PP_ALIGN.CENTER
    text(s, cxs[i] + Inches(0.12), ty, cw - Inches(0.24), Inches(0.5),
         [[(h, 11.5, WHITE, True)]], align=al, anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("Sales", "Named licensed user", DEFS["NamedLicensedUser"], "Ratified", GREEN),
    ("Finance", "Paid licensed seat", DEFS["PaidLicensedSeat"], "Ratified", GREEN),
    ("Global Service & Delivery", "Served population", DEFS["ServedPopulation"], "Proposed", AMBER),
    ("Shared Services", "Login-capable identity", DEFS["LoginCapableIdentity"], "Proposed", AMBER),
]
for r, (bu, df, ct, st, stc) in enumerate(rows):
    ry = ty + Inches(0.5) + rowh * r
    rect(s, tx, ry, tw, rowh, fill=WHITE if r % 2 == 0 else RGBColor(0xF7,0xF9,0xFB),
         line=LINE, line_w=0.5)
    text(s, cxs[0] + Inches(0.12), ry, cols[0] - Inches(0.24), rowh,
         [[(bu, 11.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cxs[1] + Inches(0.12), ry, cols[1] - Inches(0.24), rowh,
         [[(df, 11, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cxs[2] + Inches(0.12), ry, cols[2] - Inches(0.24), rowh,
         [[(f"{ct:,}", 14, NAVY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cxs[3] + Inches(0.12), ry, cols[3] - Inches(0.24), rowh,
         [[(st, 9.5, stc, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
spread = max(DEFS.values()) - min(DEFS.values())
rx = Inches(8.35); rw = Inches(4.4)
rect(s, rx, Inches(3.4), rw, Inches(1.35), fill=NAVY, round_=True)
text(s, rx + Inches(0.3), Inches(3.55), rw - Inches(0.6), Inches(0.5),
     [[("DEFINITION SPREAD (MAX \u2212 MIN)", 10, CFE0, True)]])
text(s, rx + Inches(0.3), Inches(3.9), rw - Inches(0.6), Inches(0.8),
     [[(f"{spread:,}", 40, WHITE, True), ("   across 5,000 distinct users", 12, CFE0, False)]],
     anchor=MSO_ANCHOR.MIDDLE)
rect(s, rx, Inches(4.95), rw, Inches(1.5), fill=WHITE, line=LINE, line_w=1.0, round_=True)
rect(s, rx, Inches(4.95), Inches(0.09), Inches(1.5), fill=BLUE)
text(s, rx + Inches(0.32), Inches(5.12), rw - Inches(0.6), Inches(1.2),
     [[("The agent\u2019s answer:", 11, GREY, True)],
      [("\u201CThe variance is a governance / definition issue, not a data-quality problem.\u201D",
        13, NAVY, True, "i")]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=4)
text(s, Inches(0.55), Inches(6.12), Inches(7.5), Inches(0.7),
     [[("Every answer is grounded in the trusted ", 11.5, GREY, False),
       ("ConferencesData", 11.5, NAVY, True),
       (" model \u2014 each definition states its owner and source, and reconciles from the same tables.",
        11.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# ---------------- SLIDE 2 : the foundation ----------------------------------
s = prs.slides.add_slide(BLANK); bg(s); header(s); footer(s)
text(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.6),
     [[("The trusted foundation behind that answer", 27, NAVY, True)]])
text(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.5),
     [[("A governed model, an ontology, and a conversational agent \u2014 built on Microsoft Fabric.",
        16, GREY, False)]])
lx, lw = Inches(0.55), Inches(6.0)
rect(s, lx, Inches(2.5), lw, Inches(4.05), fill=WHITE, line=LINE, line_w=1.0, round_=True)
text(s, lx + Inches(0.3), Inches(2.68), lw - Inches(0.6), Inches(0.4),
     [[("THE ONTOLOGY", 12, BLUE, True), ("   \u00b7  7 core objects, 7 verbs", 11, GREY, False)]])
verbs = [("BusinessUnit", "supports", "User"), ("BusinessUnit", "supports", "Product"),
         ("User", "has", "Role"), ("User", "has", "Licence"), ("User", "uses", "Product"),
         ("Licence", "belongs to", "Org"), ("Product", "has", "Entitlement")]
vy = Inches(3.2)
for a, v, b in verbs:
    text(s, lx + Inches(0.45), vy, lw - Inches(0.8), Inches(0.4),
         [[(a, 12.5, NAVY, True), ("  \u2014 " + v + " \u2192  ", 11, GREY, False), (b, 12.5, NAVY, True)]],
         anchor=MSO_ANCHOR.MIDDLE)
    vy = vy + Inches(0.44)
text(s, lx + Inches(0.3), Inches(6.16), lw - Inches(0.6), Inches(0.35),
     [[("Entities: ", 10.5, GREY, True),
       ("User \u00b7 Role \u00b7 BusinessUnit \u00b7 Licence \u00b7 Product \u00b7 Org \u00b7 Entitlement",
        10.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
rx, rw = Inches(6.95), Inches(5.83)
rect(s, rx, Inches(2.5), rw, Inches(4.05), fill=WHITE, line=LINE, line_w=1.0, round_=True)
text(s, rx + Inches(0.3), Inches(2.68), rw - Inches(0.6), Inches(0.4),
     [[("THE STACK", 12, BLUE, True), ("   \u00b7  signal \u2192 trusted answer", 11, GREY, False)]])
steps = [("Fabric Lakehouse", "ConferencesData \u2014 17 tables", NAVY),
         ("Semantic model (Direct Lake)", "Live measures; definition spread self-evident", NAVY),
         ("Fabric Ontology", "13 entity types, verbs bound to OneLake", NAVY),
         ("Fabric Data Agent", "Grounded, conversational, cites its definitions", BLUE),
         ("Foundry Agent + Foundry IQ", "Adds knowledge; delivered in Teams / M365 Copilot", BLUED)]
sy = Inches(3.18); bxh = Inches(0.6)
for i, (t, d, col) in enumerate(steps):
    rect(s, rx + Inches(0.35), sy, rw - Inches(0.7), bxh, fill=col, round_=True)
    text(s, rx + Inches(0.6), sy, rw - Inches(1.2), bxh,
         [[(t, 12.5, WHITE, True), ("   \u2014 " + d, 10, RGBColor(0xDB,0xE7,0xF4), False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, rx + rw/2 - Inches(0.09),
                               sy + bxh + Inches(0.005), Inches(0.18), Inches(0.12))
        a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xB6,0xC4,0xD4)
        _noline(a); a.shadow.inherit = False
    sy = sy + bxh + Inches(0.135)

# ---------------- SLIDE 3 : two questions answered --------------------------
s = prs.slides.add_slide(BLANK); bg(s); header(s); footer(s)
text(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.6),
     [[("Two questions the business actually asks", 27, NAVY, True)]])
text(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.5),
     [[("Answered live by the agent \u2014 every figure grounded in the trusted ", 16, GREY, False),
       ("ConferencesData", 16, BLUE, True), (" model.", 16, GREY, False)]])


def result_card(x, tag, q, rows):
    cw = Inches(5.95); cy = Inches(2.52); ch = Inches(4.15)
    rect(s, x, cy, cw, ch, fill=WHITE, line=LINE, line_w=1.0, round_=True)
    rect(s, x, cy + Inches(0.14), Inches(0.09), ch - Inches(0.28), fill=BLUE)
    text(s, x + Inches(0.34), cy + Inches(0.2), cw - Inches(0.6), Inches(0.35),
         [[(tag, 12, BLUE, True)]])
    rect(s, x + Inches(0.34), cy + Inches(0.62), cw - Inches(0.66), Inches(0.62),
         fill=LBLUE, round_=True)
    text(s, x + Inches(0.52), cy + Inches(0.62), cw - Inches(1.0), Inches(0.62),
         [[("Q   ", 12, BLUE, True), (q, 11.5, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    ry = cy + Inches(1.46); rh = Inches(0.44)
    for i, (nm, val) in enumerate(rows):
        if i == 0:
            rect(s, x + Inches(0.34), ry, cw - Inches(0.66), rh, fill=RGBColor(0xEF,0xF5,0xFC))
        text(s, x + Inches(0.52), ry, Inches(3.2), rh,
             [[(f"{i+1}.  ", 11, GREY, False), (nm, 13.5 if i == 0 else 12, NAVY if i == 0 else INK, True)]],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + Inches(3.5), ry, cw - Inches(3.5) - Inches(0.4), rh,
             [[(val, 13 if i == 0 else 11, NAVY if i == 0 else GREY, i == 0)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        ry = ry + rh
    text(s, x + Inches(0.34), cy + ch - Inches(0.4), cw - Inches(0.6), Inches(0.32),
         [[("\u25CF  via Fabric Data Agent", 9.5, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)


result_card(Inches(0.55), "SPEAKER IMPACT  \u00b7  licensed users attended",
            "Which speakers had the most licensed users attend?",
            [(n, f"{c:,}") for n, c in TOP_SPK[:5]])
result_card(Inches(6.83), "SPONSOR VALUE  \u00b7  influenced pipeline",
            "What sponsors drove the most value?",
            [(n, (f"{fmt_m(p)}  \u00b7  {r}\u00d7" if i == 0 else fmt_m(p)))
             for i, (n, p, r) in enumerate(TOP_SPN[:5])])

# ---------------- SLIDE 4 : one agent, two grounded brains ------------------
s = prs.slides.add_slide(BLANK); bg(s); header(s); footer(s)
text(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.6),
     [[("One agent, two grounded brains", 27, NAVY, True)]])
text(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.5),
     [[("The Foundry orchestrator routes each question to the right source \u2014 and tells you which one it used.",
        16, GREY, False)]])
ob_w = Inches(6.4); ob_x = (SW - ob_w) / 2; ob_y = Inches(2.5)
rect(s, ob_x, ob_y, ob_w, Inches(0.92), fill=NAVY, round_=True)
text(s, ob_x + Inches(0.3), ob_y + Inches(0.13), ob_w - Inches(0.6), Inches(0.36),
     [[("Foundry Agent  \u00b7  ConferenceIQ Orchestrator", 14, WHITE, True), ("   (gpt-4o)", 11, CFE0, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, ob_x + Inches(0.3), ob_y + Inches(0.5), ob_w - Inches(0.6), Inches(0.34),
     [[("Reads the question, picks the source, cites what it used", 10.5, CFE0, False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
col_w = Inches(5.75); col_y = Inches(3.95); col_h = Inches(2.05)
lx = Inches(0.65); rxx = SW - Inches(0.65) - col_w
tspk = f"{TOP_SPK[0][0]} {TOP_SPK[0][1]:,}"
tspn = f"{TOP_SPN[0][0]} {fmt_m(TOP_SPN[0][1])}"
for cx, htitle, subttl, bullets, ex in [
    (lx, "FABRIC DATA AGENT", "Source of truth for numbers",
     "counts  \u00b7  rankings  \u00b7  ROI  \u00b7  influenced pipeline  \u00b7  attendance",
     f"e.g.  {tspk}   \u00b7   {tspn}"),
    (rxx, "FOUNDRY IQ KNOWLEDGE", "Qualitative context & narrative",
     "speaker bios  \u00b7  session guides  \u00b7  sponsor prospectus  \u00b7  FAQ  \u00b7  policies",
     "e.g.  Platinum package $250K   \u00b7   speaker backgrounds"),
]:
    rect(s, cx, col_y, col_w, col_h, fill=WHITE, line=LINE, line_w=1.0, round_=True)
    rect(s, cx, col_y, col_w, Inches(0.07), fill=BLUE)
    text(s, cx + Inches(0.32), col_y + Inches(0.2), col_w - Inches(0.6), Inches(0.35),
         [[(htitle, 13, NAVY, True)]])
    text(s, cx + Inches(0.32), col_y + Inches(0.58), col_w - Inches(0.6), Inches(0.32),
         [[(subttl, 11.5, BLUE, True)]])
    text(s, cx + Inches(0.32), col_y + Inches(0.98), col_w - Inches(0.6), Inches(0.4),
         [[(bullets, 11, GREY, False)]], line_spacing=1.05)
    text(s, cx + Inches(0.32), col_y + Inches(1.5), col_w - Inches(0.6), Inches(0.42),
         [[(ex, 10.5, INK, False, "i")]], anchor=MSO_ANCHOR.MIDDLE)
for ax in [lx + col_w / 2, rxx + col_w / 2]:
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, ax - Inches(0.1),
                           ob_y + Inches(0.92) + Inches(0.06), Inches(0.2), Inches(0.4))
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xB6,0xC4,0xD4)
    _noline(a); a.shadow.inherit = False
rect(s, Inches(0.65), Inches(6.28), SW - Inches(1.3), Inches(0.5), fill=LBLUE, round_=True)
text(s, Inches(0.95), Inches(6.28), SW - Inches(1.9), Inches(0.5),
     [[("Blended questions use both \u2014 ", 11.5, NAVY, True),
       ("a hard number from Fabric, enriched with knowledge \u2014 and every answer names its source.",
        11.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ---------------- SLIDE 5 : ask it anything ---------------------------------
s = prs.slides.add_slide(BLANK); bg(s); header(s); footer(s)
text(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.6),
     [[("Ask it anything \u2014 even questions that need both", 27, NAVY, True)]])
text(s, Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.5),
     [[("Blended questions force the agent to combine a hard number with descriptive context.",
        16, GREY, False)]])
blended = [
    ("Which speaker drew the most licensed users \u2014 and what\u2019s their background?",
     f"\u2192  {TOP_SPK[0][0]} ({TOP_SPK[0][1]:,}), enriched with their speaker bio"),
    ("Our top sponsor by influenced pipeline \u2014 which package did they buy?",
     f"\u2192  {TOP_SPN[0][0]} ({fmt_m(TOP_SPN[0][1])}), with their sponsorship tier from the prospectus"),
    ("Which conference had the highest licensed attendance, and what were its themes?",
     "\u2192  Top event by registrations, plus its agenda and headline tracks"),
    ("Top sponsors by ROI \u2014 are they on our premium tiers, and what would an upgrade add?",
     "\u2192  ROI ranking, matched to package benefits and upsell options"),
]
cy = Inches(2.5); cardh = Inches(0.82); gap = Inches(0.12)
for q, prev in blended:
    rect(s, Inches(0.55), cy, Inches(12.23), cardh, fill=WHITE, line=LINE, line_w=1.0, round_=True)
    rect(s, Inches(0.55), cy + Inches(0.1), Inches(0.09), cardh - Inches(0.2), fill=BLUE)
    text(s, Inches(0.85), cy + Inches(0.11), Inches(9.1), Inches(0.36),
         [[(q, 13.5, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.85), cy + Inches(0.45), Inches(9.1), Inches(0.3),
         [[(prev, 11, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(10.2), cy + Inches(0.23), Inches(2.3), Inches(0.36), fill=LBLUE, round_=True)
    text(s, Inches(10.2), cy + Inches(0.23), Inches(2.3), Inches(0.36),
         [[("Fabric + Foundry IQ", 10, BLUED, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cy = cy + cardh + gap
text(s, Inches(0.55), Inches(6.4), Inches(12.2), Inches(0.42),
     [[("Delivered where the business already works \u2014 ", 12, GREY, False),
       ("Teams, Microsoft 365 Copilot, or a branded web client.", 12, NAVY, True)]],
     anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("saved", OUT)
